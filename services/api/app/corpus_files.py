import csv
import io
import json
import re
from functools import lru_cache
from html.parser import HTMLParser
from pathlib import Path

from docx import Document
from markitdown import MarkItDown
from openpyxl import load_workbook
from pptx import Presentation
from pypdf import PdfReader

SUPPORTED_CORPUS_EXTENSIONS = {
    ".txt",
    ".md",
    ".markdown",
    ".csv",
    ".json",
    ".html",
    ".htm",
    ".pdf",
    ".docx",
    ".xlsx",
    ".pptx",
}
MAX_FILE_BYTES = 10 * 1024 * 1024
MAX_EXTRACTED_CHARS = 120_000
CORPUS_CHUNK_CHARS = 11_500
MARKITDOWN_EXTENSIONS = {".html", ".htm", ".pdf", ".docx", ".pptx"}


class CorpusFileError(ValueError):
    pass


class _TextHTMLParser(HTMLParser):
    def __init__(self) -> None:
        super().__init__()
        self.parts: list[str] = []
        self._ignored_depth = 0

    def handle_starttag(self, tag: str, attrs: list[tuple[str, str | None]]) -> None:
        if tag in {"script", "style"}:
            self._ignored_depth += 1
        elif tag in {"br", "p", "div", "li", "tr", "h1", "h2", "h3", "h4"}:
            self.parts.append("\n")

    def handle_endtag(self, tag: str) -> None:
        if tag in {"script", "style"}:
            self._ignored_depth = max(0, self._ignored_depth - 1)
        elif tag in {"p", "div", "li", "tr"}:
            self.parts.append("\n")

    def handle_data(self, data: str) -> None:
        if not self._ignored_depth:
            self.parts.append(data)


def parse_corpus_file(filename: str, data: bytes) -> str:
    extension = Path(filename).suffix.lower()
    if extension not in SUPPORTED_CORPUS_EXTENSIONS:
        raise CorpusFileError(f"不支持 {extension or '无扩展名'} 格式")
    if not data:
        raise CorpusFileError("文件内容为空")
    if len(data) > MAX_FILE_BYTES:
        raise CorpusFileError("文件超过 10 MB")

    try:
        if extension in {".txt", ".md", ".markdown"}:
            text = _decode_text(data)
        elif extension == ".csv":
            text = _parse_csv(data)
        elif extension == ".json":
            text = _parse_json(data)
        elif extension in MARKITDOWN_EXTENSIONS:
            text = _parse_rich_document(extension, data)
        elif extension == ".xlsx":
            text = _parse_xlsx(data)
    except CorpusFileError:
        raise
    except Exception as exc:
        raise CorpusFileError("文件损坏、加密或无法解析") from exc

    text = _normalize_text(text)
    if not text:
        raise CorpusFileError("未提取到可用文字")
    if len(text) > MAX_EXTRACTED_CHARS:
        raise CorpusFileError("提取文字超过 12 万字，请拆分文件后重试")
    return text


def split_corpus_text(text: str, limit: int = CORPUS_CHUNK_CHARS) -> list[str]:
    if len(text) <= limit:
        return [text]

    chunks: list[str] = []
    remaining = text
    while remaining:
        if len(remaining) <= limit:
            chunks.append(remaining.strip())
            break
        boundary = max(
            remaining.rfind("\n\n", 0, limit + 1),
            remaining.rfind("\n", 0, limit + 1),
            remaining.rfind("。", 0, limit + 1),
            remaining.rfind("；", 0, limit + 1),
        )
        if boundary < limit // 2:
            boundary = limit
        elif remaining[boundary : boundary + 1] in {"。", "；"}:
            boundary += 1
        chunk = remaining[:boundary].strip()
        if chunk:
            chunks.append(chunk)
        remaining = remaining[boundary:].strip()
    return chunks


def _decode_text(data: bytes) -> str:
    for encoding in ("utf-8-sig", "gb18030", "utf-16"):
        try:
            return data.decode(encoding)
        except UnicodeDecodeError:
            continue
    raise CorpusFileError("文本编码无法识别，请转为 UTF-8 后重试")


def _parse_csv(data: bytes) -> str:
    decoded = _decode_text(data)
    sample = decoded[:4096]
    try:
        dialect = csv.Sniffer().sniff(sample, delimiters=",\t;|")
    except csv.Error:
        dialect = csv.excel
    rows = csv.reader(io.StringIO(decoded), dialect)
    return "\n".join("\t".join(cell.strip() for cell in row) for row in rows)


def _parse_json(data: bytes) -> str:
    try:
        value = json.loads(_decode_text(data))
    except json.JSONDecodeError as exc:
        raise CorpusFileError("JSON 格式无效") from exc
    return json.dumps(value, ensure_ascii=False, indent=2)


def _parse_html(data: bytes) -> str:
    parser = _TextHTMLParser()
    parser.feed(_decode_text(data))
    return "".join(parser.parts)


@lru_cache(maxsize=1)
def _get_markitdown() -> MarkItDown:
    return MarkItDown(enable_plugins=False)


def _parse_rich_document(extension: str, data: bytes) -> str:
    """Convert rich documents to Markdown, with the legacy parser as fallback."""
    try:
        result = _get_markitdown().convert_stream(
            io.BytesIO(data),
            file_extension=extension,
        )
        markdown = result.markdown.strip()
        if markdown:
            return markdown
    except Exception:
        # The established parsers remain available for damaged or unusual files.
        pass

    if extension in {".html", ".htm"}:
        return _parse_html(data)
    if extension == ".pdf":
        return _parse_pdf(data)
    if extension == ".docx":
        return _parse_docx(data)
    return _parse_pptx(data)


def _parse_pdf(data: bytes) -> str:
    reader = PdfReader(io.BytesIO(data))
    if reader.is_encrypted:
        try:
            reader.decrypt("")
        except Exception as exc:
            raise CorpusFileError("暂不支持加密 PDF") from exc
    return "\n\n".join(page.extract_text() or "" for page in reader.pages)


def _parse_docx(data: bytes) -> str:
    document = Document(io.BytesIO(data))
    parts = [paragraph.text for paragraph in document.paragraphs]
    for table in document.tables:
        for row in table.rows:
            parts.append("\t".join(cell.text.strip() for cell in row.cells))
    return "\n".join(parts)


def _parse_xlsx(data: bytes) -> str:
    workbook = load_workbook(io.BytesIO(data), read_only=True, data_only=True)
    parts: list[str] = []
    try:
        for sheet in workbook.worksheets:
            parts.append(f"## 工作表：{sheet.title}")
            for row in sheet.iter_rows(values_only=True):
                values = ["" if value is None else str(value).strip() for value in row]
                if any(values):
                    parts.append("\t".join(values).rstrip())
                if sum(len(part) for part in parts) > MAX_EXTRACTED_CHARS:
                    raise CorpusFileError("提取文字超过 12 万字，请拆分文件后重试")
    finally:
        workbook.close()
    return "\n".join(parts)


def _parse_pptx(data: bytes) -> str:
    presentation = Presentation(io.BytesIO(data))
    parts: list[str] = []
    for index, slide in enumerate(presentation.slides, start=1):
        parts.append(f"## 第 {index} 页")
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                parts.append(shape.text.strip())
    return "\n".join(parts)


def _normalize_text(text: str) -> str:
    text = text.replace("\x00", "").replace("\r\n", "\n").replace("\r", "\n")
    text = re.sub(r"[ \t]+\n", "\n", text)
    text = re.sub(r"\n{3,}", "\n\n", text)
    return text.strip()
