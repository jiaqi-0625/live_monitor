import io
import json
from pathlib import Path

from docx import Document
from fastapi.testclient import TestClient
from openpyxl import Workbook
from pptx import Presentation
from pydantic import SecretStr

from app import ai_analysis as ai_analysis_module
from app import main as main_module
from app.config import Settings
from app.corpus_files import CorpusFileError, parse_corpus_file, split_corpus_text
from app.database import Database


def test_parse_text_csv_json_docx_and_xlsx():
    assert parse_corpus_file("口径.txt", "中文卖点".encode()) == "中文卖点"
    assert "车型\t卖点" in parse_corpus_file("资料.csv", "车型,卖点\nA,安全".encode())
    parsed_json = parse_corpus_file(
        "规则.json",
        json.dumps({"阈值": 20}, ensure_ascii=False).encode(),
    )
    assert '"阈值": 20' in parsed_json

    document = Document()
    document.add_heading("Word 车型卖点", level=1)
    document.add_paragraph("全系标配主动安全系统")
    docx_buffer = io.BytesIO()
    document.save(docx_buffer)
    parsed_docx = parse_corpus_file("车型.docx", docx_buffer.getvalue())
    assert "# Word 车型卖点" in parsed_docx
    assert "全系标配主动安全系统" in parsed_docx

    workbook = Workbook()
    sheet = workbook.active
    sheet.title = "价格政策"
    sheet.append(["版本", "优惠"])
    sheet.append(["旗舰版", 20000])
    xlsx_buffer = io.BytesIO()
    workbook.save(xlsx_buffer)
    parsed_xlsx = parse_corpus_file("政策.xlsx", xlsx_buffer.getvalue())
    assert "工作表：价格政策" in parsed_xlsx
    assert "旗舰版\t20000" in parsed_xlsx

    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[5])
    slide.shapes.title.text = "直播流量层级"
    pptx_buffer = io.BytesIO()
    presentation.save(pptx_buffer)
    assert "直播流量层级" in parse_corpus_file("流量.pptx", pptx_buffer.getvalue())

    parsed_html = parse_corpus_file(
        "人群.html",
        "<h1>目标人群</h1><table><tr><td>高意向</td></tr></table>".encode(),
    )
    assert "# 目标人群" in parsed_html
    assert "高意向" in parsed_html


def test_parse_rejects_unsupported_and_splits_long_text():
    try:
        parse_corpus_file("图片.png", b"not-an-image")
    except CorpusFileError as exc:
        assert "不支持" in str(exc)
    else:
        raise AssertionError("unsupported format should fail")

    chunks = split_corpus_text(("一段内容。\n\n" * 3000).strip(), limit=1000)
    assert len(chunks) > 1
    assert all(len(chunk) <= 1000 for chunk in chunks)


def test_corpus_import_api_handles_partial_failure(tmp_path: Path, monkeypatch):
    test_database = Database(tmp_path / "corpus-import.db")
    monkeypatch.setattr(main_module, "database", test_database)
    monkeypatch.setattr(main_module.settings, "llm_api_key", SecretStr("test-key"))
    monkeypatch.setattr(
        main_module,
        "condense_corpus_markdown",
        lambda settings, *, filename, category, source_text: (
            f"# {Path(filename).stem}\n> 分类：车型卖点｜来源：{filename}\n\n"
            f"## 核心摘要\n- {source_text}"
        ),
    )

    with TestClient(main_module.app) as client:
        response = client.post(
            "/api/corpus/import",
            data={"operator_name": "优化师甲", "category": "vehicle"},
            files=[
                ("files", ("星越L.txt", "全系标配主动安全系统。".encode(), "text/plain")),
                ("files", ("海报.png", b"invalid", "image/png")),
            ],
        )

    assert response.status_code == 200
    result = response.json()
    assert result["imported_files"] == 1
    assert result["imported_entries"] == 1
    assert result["entries"][0]["source_name"] == "星越L.txt"
    assert result["entries"][0]["content"].startswith("# 星越L")
    assert result["original_chars"] > 0
    assert result["saved_chars"] > 0
    assert result["fallback_files"] == []
    assert result["failures"][0]["filename"] == "海报.png"


def test_corpus_import_requires_llm(tmp_path: Path, monkeypatch):
    monkeypatch.setattr(main_module, "database", Database(tmp_path / "no-llm.db"))
    monkeypatch.setattr(main_module.settings, "llm_api_key", None)

    with TestClient(main_module.app) as client:
        response = client.post(
            "/api/corpus/import",
            data={"operator_name": "优化师甲", "category": "vehicle"},
            files={"files": ("车型.txt", b"content", "text/plain")},
        )

    assert response.status_code == 503
    assert "大模型服务尚未配置" in response.json()["detail"]


def test_corpus_ai_timeout_has_actionable_message(monkeypatch):
    def raise_timeout(*args, **kwargs):
        raise TimeoutError

    monkeypatch.setattr(ai_analysis_module, "urlopen", raise_timeout)
    configured_settings = Settings(llm_api_key=SecretStr("test-key"))

    try:
        ai_analysis_module._call_llm_text(
            configured_settings,
            system_prompt="system",
            user_prompt="content",
        )
    except RuntimeError as exc:
        assert str(exc) == "大模型处理超时，请稍后重试"
    else:
        raise AssertionError("model timeout should have an actionable error")


def test_normalize_corpus_markdown_accepts_explanation_and_fence():
    normalized = ai_analysis_module._normalize_corpus_markdown_output(
        """下面是整理结果：
```Markdown
# 流量层级和人群问题

## 核心摘要
- 直播间流量层级需要结合原文判断。
```
""",
        filename="直播间账号流量层级和人群的问题 .pdf",
        category_label="其他资料",
    )

    assert normalized.startswith("# 流量层级和人群问题")
    assert "来源：直播间账号流量层级和人群的问题 .pdf" in normalized
    assert "## 关键事实与口径\n无" in normalized
    assert "```" not in normalized


def test_normalize_corpus_markdown_adds_missing_title_and_sections():
    normalized = ai_analysis_module._normalize_corpus_markdown_output(
        "## 核心摘要\n- 人群标签存在偏差。",
        filename="账号人群.pdf",
        category_label="流量与观看",
    )

    assert normalized.startswith("# 账号人群\n> 分类：流量与观看｜来源：账号人群.pdf")
    assert "## 核心摘要\n- 人群标签存在偏差。" in normalized
    assert "## 可执行话术与动作\n无" in normalized


def test_normalize_corpus_markdown_wraps_plain_text_and_rejects_empty():
    normalized = ai_analysis_module._normalize_corpus_markdown_output(
        "- 流量层级为测试结论。",
        filename="流量层级.pdf",
        category_label="其他资料",
    )
    assert "## 核心摘要\n- 流量层级为测试结论。" in normalized

    try:
        ai_analysis_module._normalize_corpus_markdown_output(
            "  ",
            filename="空文件.pdf",
            category_label="其他资料",
        )
    except RuntimeError as exc:
        assert str(exc) == "大模型未返回可用语料"
    else:
        raise AssertionError("empty model output should fail")


def test_call_llm_text_rejects_empty_visible_content(monkeypatch):
    class FakeResponse:
        def __enter__(self):
            return self

        def __exit__(self, exc_type, exc, traceback):
            return False

        def read(self):
            return json.dumps(
                {
                    "choices": [
                        {
                            "message": {
                                "content": "",
                                "reasoning_content": "内部推理内容",
                            },
                            "finish_reason": "length",
                        }
                    ]
                }
            ).encode()

    monkeypatch.setattr(ai_analysis_module, "urlopen", lambda *args, **kwargs: FakeResponse())

    try:
        ai_analysis_module._call_llm_text(
            Settings(llm_api_key=SecretStr("test-key")),
            system_prompt="system",
            user_prompt="content",
        )
    except ai_analysis_module.EmptyLlmContentError as exc:
        assert "length" in str(exc)
    else:
        raise AssertionError("empty model content should be detected")


def test_condense_corpus_falls_back_to_extracted_source(monkeypatch):
    monkeypatch.setattr(
        ai_analysis_module,
        "_call_llm_text",
        lambda *args, **kwargs: (_ for _ in ()).throw(
            ai_analysis_module.EmptyLlmContentError("empty")
        ),
    )

    markdown = ai_analysis_module.condense_corpus_markdown(
        Settings(llm_api_key=SecretStr("test-key")),
        filename="直播间账号流量层级和人群的问题 .pdf",
        category="other",
        source_text="账号流量层级与目标人群存在偏差。",
    )

    assert markdown.startswith("# 直播间账号流量层级和人群的问题")
    assert "已按原文导入" in markdown
    assert "账号流量层级与目标人群存在偏差。" in markdown


def test_v4_corpus_request_disables_thinking(monkeypatch):
    captured_payload = None

    class FakeResponse:
        def __enter__(self):
            return self

        def __exit__(self, exc_type, exc, traceback):
            return False

        def read(self):
            return json.dumps(
                {
                    "choices": [
                        {
                            "message": {"content": "# 正文"},
                            "finish_reason": "stop",
                        }
                    ]
                }
            ).encode()

    def capture_request(request, **kwargs):
        nonlocal captured_payload
        captured_payload = json.loads(request.data.decode())
        return FakeResponse()

    monkeypatch.setattr(ai_analysis_module, "urlopen", capture_request)
    ai_analysis_module._call_llm_text(
        Settings(llm_api_key=SecretStr("test-key"), llm_model="deepseek-v4-flash"),
        system_prompt="system",
        user_prompt="content",
    )

    assert captured_payload["thinking"] == {"type": "disabled"}
